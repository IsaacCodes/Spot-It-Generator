#USER DEFINED CONSTANTS
CARD_IMAGE_DIRECTORY = "images"
PROJECTIVE_PLANE_SHEET = "projective_plane.xlsx"
SPOT_IT_GAME_SLIDESHOW = "spot_it_game.pptx"


#DO NOT MODIFY CODE BELOW HERE

#Imports and set up
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from PIL import Image
import os
from sys import exit
import random

#Loads in the images from the directory
images_path = os.path.join(os.path.curdir, CARD_IMAGE_DIRECTORY)
images_names = os.listdir(images_path)
images_path_names = [os.path.join(images_path, image_name) for image_name in images_names]

#Checks for right file extensions and count
valid_extensions = ["png", "jpg", "jpeg", "gif"]
for image in images_path_names:
  if not os.path.isfile(image):
    print("Please do not place any subfolders in the images folder")
    exit()
  extension = image.split(os.extsep)[-1]
  if extension not in valid_extensions:
    print(f"File type '{extension}' is not supported. Please use one of the following: {', '.join(valid_extensions)}")
    exit()

if len(images_names) != 57:
  print(f"You have placed '{len(images_names)}' files in images. Please instead place 57.")
  exit()

#Loads in the card indicies from the projective plane
cards_df = pd.read_excel(PROJECTIVE_PLANE_SHEET)
cards = cards_df.values.tolist()
random.shuffle(cards)
print(cards_df)

#Grid position for slide
grid_positions = [(-1, -1), (0, -1), (1, -1), (-1, 0), (1, 0), (-1, 1), (0, 1), (1, 1)]

#Slide and image dimensions
slide_width = Inches(10)
slide_height = Inches(10)
img_width = Inches(2.5)
img_height = Inches(2.5)
x_spacing = Inches(3)
y_spacing = Inches(2.5)

#Create PowerPoint
pres = Presentation()
blank_layout = pres.slide_layouts[6]

for card in cards:
  random.shuffle(grid_positions)
  slide = pres.slides.add_slide(blank_layout)

  for i, symbol_num in enumerate(card[1:]):
    col, row = grid_positions[i]

    img_path = images_path_names[symbol_num-1]

    img = Image.open(img_path)

    if img.width > img.height:
      slide_img = slide.shapes.add_picture(img_path, 0, 0, width=img_width * random.uniform(0.7, 1))
    else:
      slide_img = slide.shapes.add_picture(img_path, 0, 0, height=img_height * random.uniform(0.7, 1))

    slide_img.left = int(slide_width / 2 + col * x_spacing - slide_img.width / 2)
    slide_img.top = int(slide_height / 2 + row * y_spacing - slide_img.height / 2)

    slide_img.rotation = random.randint(0, 360)

    print(img_path)

#Save
pres.save(SPOT_IT_GAME_SLIDESHOW)
print(f"Presentation successfully saved as {SPOT_IT_GAME_SLIDESHOW}")